#!/usr/bin/env python3
"""Bouw een Cartel-deck uit een slideplan, op het sjabloon.

Het slideplan is een json-bestand: per slide de layout en de tekst per rol. Dit script
zet dat om naar een pptx door slides op het sjabloon aan te maken en de placeholders te
vullen volgens hun gemeten rol. Het ontwerpt niets en het verzint geen posities.

    python3 bouw-deck.py plan.json -o deck.pptx

Vorm van het plan:

{
  "sjabloon": "assets/cartel-basis.pptx",
  "titel": "Novo Nordisk, obesitas, research",
  "slides": [
    {
      "layout": "CartelAgency-Titeldia-4",
      "action title": "Novo Nordisk: obesitas in België",
      "tekst": "Researchdossier, augustus 2026",
      "notities": "Kort kaderen, niet voorlezen."
    },
    {
      "layout": "CartelAgency-Titel+Tekst-Zwart",
      "action title": "Alleen diabetes wordt terugbetaald, obesitas zelf niet",
      "tekst": "Mounjaro is het enige middel met terugbetaling, en enkel bij diabetes.",
      "bron": "Bron: RIZIV, 2026",
      "notities": "Dit is de kern van het marktblok."
    },
    {
      "layout": "CartelAgency-Tussentitel-2",
      "action title": "De markt",
      "kort label": "01",
      "beeld": "sfeerbeeld van een apotheek, breed formaat"
    }
  ]
}

Rollen die je kan invullen: `action title`, `tekst`, `kort label`, `bron`, `extra`,
`beeld` (omschrijving voor de plaatshouder) en `notities`. Welke rollen een layout heeft,
staat in `assets/layouts.md`.
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from capaciteit import rollen, TOLERANTIE  # noqa: E402

INVULBAAR = ("action title", "tekst", "kort label", "bron", "extra")


def zoek_layout(prs, naam):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == naam:
                return layout, master
    # tweede poging: zonder master-prefix vergelijken
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.split("_", 1)[-1] == naam.split("_", 1)[-1]:
                return layout, master
    return None, None


def zet_tekst(ph, tekst):
    """Vervang de inhoud van een placeholder zonder de opmaak te slopen.

    Niet `text_frame.text = ...` gebruiken: dat plet de paragraaf tot één run zonder
    opmaak. Per regel een paragraaf, en de opmaak van de eerste paragraaf hergebruiken.
    """
    kader = ph.text_frame
    regels = [r for r in str(tekst).split("\n")]
    eerste = kader.paragraphs[0]
    for overtollig in list(eerste.runs):
        overtollig._r.getparent().remove(overtollig._r)
    eerste.add_run().text = regels[0]
    for regel in regels[1:]:
        para = kader.add_paragraph()
        para.add_run().text = regel


def zet_beeld(slide, ph, omschrijving):
    """Vervang een beeldplaceholder door een rechthoek met omschrijving."""
    links, boven, breed, hoog = ph.left, ph.top, ph.width, ph.height
    ph._element.getparent().remove(ph._element)
    vorm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, links, boven, breed, hoog)
    vorm.fill.solid()
    vorm.fill.fore_color.rgb = RGBColor(0x9F, 0x9F, 0x9F)
    vorm.line.color.rgb = RGBColor(0x42, 0x42, 0x42)
    vorm.line.width = Pt(1)
    kader = vorm.text_frame
    kader.word_wrap = True
    kader.text = f"PLAATSHOUDER BEELD\n{omschrijving}"
    for para in kader.paragraphs:
        for run in para.runs:
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0x14, 0x14, 0x14)
    return vorm


def bouw(plan, uit, sjabloon=None, streng=False):
    pad = sjabloon or plan.get("sjabloon")
    if not pad or not os.path.exists(pad):
        raise SystemExit(f"Sjabloon niet gevonden: {pad}")
    prs = Presentation(pad)

    # het sjabloon bevat voorbeeldslides: die gaan eruit
    lijst = prs.slides._sldIdLst
    from pptx.oxml.ns import qn

    for sldId in list(lijst):
        prs.part.drop_rel(sldId.get(qn("r:id")))
        lijst.remove(sldId)

    waarschuwingen = []
    for nummer, opgave in enumerate(plan["slides"], start=1):
        naam = opgave.get("layout")
        layout, master = zoek_layout(prs, naam)
        if layout is None:
            waarschuwingen.append(f"slide {nummer}: layout '{naam}' bestaat niet in het sjabloon")
            continue
        slide = prs.slides.add_slide(layout)
        kaart = rollen(layout, master)

        for ph in list(slide.placeholders):
            info = kaart.get(ph.placeholder_format.idx, {})
            rol = info.get("rol", "onbekend")
            if rol == "picture":
                omschrijving = opgave.get("beeld")
                if omschrijving:
                    zet_beeld(slide, ph, omschrijving)
                continue
            if rol in ("chart", "table", "sjabloonelement", "onbekend"):
                continue
            waarde = opgave.get(rol)
            if waarde is None:
                continue
            ruimte = info.get("max_tekens")
            if ruimte and len(str(waarde)) > ruimte * TOLERANTIE:
                waarschuwingen.append(
                    f"slide {nummer} ({naam}): {rol} is {len(str(waarde))} tekens, "
                    f"er past er ongeveer {ruimte}. Kort in of kies een andere layout."
                )
                if streng:
                    continue
            zet_tekst(ph, waarde)

        gevraagd = {k for k in opgave if k in INVULBAAR}
        beschikbaar = {v.get("rol") for v in kaart.values()}
        for rol in gevraagd - beschikbaar:
            if rol == "bron" and "tekst" in beschikbaar:
                # De meeste werkslides hebben geen apart bronvak. In de echte decks staat
                # de bronregel dan onderaan de tekst. Dat gedrag nemen we over.
                idx = next(i for i, v in kaart.items() if v.get("rol") == "tekst")
                ph = slide.placeholders[idx]
                huidig = ph.text_frame.text
                zet_tekst(ph, f"{huidig}\n{opgave['bron']}" if huidig else opgave["bron"])
                continue
            waarschuwingen.append(
                f"slide {nummer} ({naam}): rol '{rol}' bestaat niet in deze layout, tekst is weggevallen"
            )

        notities = opgave.get("notities")
        if notities:
            slide.notes_slide.notes_text_frame.text = str(notities)
        else:
            waarschuwingen.append(f"slide {nummer}: geen spreeknotities")

    prs.save(uit)
    return waarschuwingen


def main():
    arg = argparse.ArgumentParser(description=__doc__)
    arg.add_argument("plan", help="pad naar het slideplan (json)")
    arg.add_argument("-o", "--uit", required=True, help="pad naar het deck")
    arg.add_argument("--sjabloon", help="overschrijft het sjabloon uit het plan")
    arg.add_argument("--streng", action="store_true",
                     help="tekst die niet past wordt weggelaten in plaats van geplaatst")
    opties = arg.parse_args()

    with open(opties.plan, encoding="utf-8") as bestand:
        plan = json.load(bestand)

    waarschuwingen = bouw(plan, opties.uit, opties.sjabloon, opties.streng)
    print(f"Gebouwd: {opties.uit} ({len(plan['slides'])} slides gevraagd)")
    if waarschuwingen:
        print(f"\n{len(waarschuwingen)} waarschuwingen:")
        for w in waarschuwingen:
            print("  -", w)
        print("\nLos deze op in het slideplan, niet in het bestand.")
    else:
        print("Geen waarschuwingen.")


if __name__ == "__main__":
    main()

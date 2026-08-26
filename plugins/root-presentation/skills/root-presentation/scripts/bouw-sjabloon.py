#!/usr/bin/env python3
"""Bouw het magere Cartel-strategiesjabloon uit een bestaand Cartel-deck.

Waarom dit script bestaat: het brede bureausjabloon heeft tientallen layouts die het
strategieteam nooit gebruikt, en het mist net de twee layouts die het het vaakst wél
gebruikt. Dit script snijdt een bestaand deck terug tot de kortlijst, gooit alle
inhoud en beeld eruit, en zet er per layout één neutrale voorbeeldslide in.

Die voorbeeldslides worden gemaakt met exact hetzelfde mechanisme als waarmee de agent
straks slides bouwt (een nieuwe slide op een layout, placeholders invullen). Ziet het
sjabloon er goed uit, dan ziet de output van de agent er ook goed uit.

Gebruik:
    python3 bouw-sjabloon.py bron.pptx -o cartel-basis.pptx
    python3 bouw-sjabloon.py bron.pptx -o cartel-basis.pptx --zonder-voorbeelden
"""

import argparse
import os
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from capaciteit import rollen  # noqa: E402

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# De kortlijst: gemeten op 1031 slides uit elf finale strategiedecks van Cartel.
# Volgorde bepaalt de volgorde van de voorbeeldslides in het sjabloon.
KORTLIJST = [
    ("Titeldia-4", "Titelslide van het deck"),
    ("Titeldia-1", "Titelslide met beeld"),
    ("Index-1", "Inhoudsopgave of overzicht van de pijlers"),
    ("Tussentitel-2", "Divider tussen twee blokken, met beeld"),
    ("Tussentitel-1", "Divider zonder beeld"),
    ("Titel+Tekst-Zwart", "De werkslide. Donkere achtergrond, witte tekst"),
    ("Titel+Tekst", "De werkslide op lichte achtergrond"),
    ("Sentence-Black", "Eén uitspraak of kernzin, donker"),
    ("Sentence", "Eén uitspraak of kernzin, licht"),
    ("Quote", "Citaat met bron"),
    ("Titel+Tekst+Foto2", "Tekst links, beeld rechts"),
    ("Titel+Tekst+Foto", "Tekst met beeld, andere verhouding"),
    ("Titel+Bullets+Foto-1-Zwart", "Bullets met beeld, donker"),
    ("Titel+Bullets+Foto-2", "Bullets met beeld, licht"),
    ("Titel+Bullets+Grafiek", "Grafiek met duiding ernaast"),
    ("Titel+Grafiek", "Eén grafiek groot"),
    ("Grafiek x2", "Twee grafieken naast elkaar"),
    ("Titel+Tekst+Titel+Tabel", "Tabel met duiding"),
    ("Tabel-fullscreen", "Tabel over de volle breedte"),
    ("Tekst-fullscreen", "Tekst over de volle breedte"),
    ("Oplijsting", "Opsomming van punten onder elkaar"),
    ("Titel+Blokken", "Drie of meer blokken naast elkaar"),
    ("Titel+Genummerde Blokken", "Genummerde stappen of principes"),
    ("Titel+Timeline", "Tijdlijn of fasering"),
    ("Blanco-zwart", "Noodgeval: lege donkere slide"),
    ("Blanco-wit", "Noodgeval: lege lichte slide"),
    ("Endslide", "Slotslide"),
]

VOORBEELDTEKST = {
    "titel": "Hier komt de action title: de conclusie, niet het onderwerp",
    "body": "Hier komt de onderbouwing. Kort, beschrijvend, zonder beeldspraak.",
    "cijfer": "58%",
    "bron": "Bron: naam van de bron, jaartal",
    "beeld": "PLAATSHOUDER BEELD\nomschrijf hier wat er komt",
}


def kern(naam):
    """Haal de master-prefix en het CartelAgency-voorvoegsel van een layoutnaam."""
    naam = re.sub(r"^\d+_", "", naam)
    return naam.replace("CartelAgency-", "").replace("CartelAgency_", "")


def kies_layouts(master):
    """Zoek per kortlijstnaam de beste layout in deze master."""
    per_kern = {}
    for layout in master.slide_layouts:
        k = kern(layout.name)
        # zonder prefix wint van met prefix: dat is de oorspronkelijke layout
        if k not in per_kern or not layout.name.startswith(tuple("0123456789")):
            per_kern.setdefault(k, layout)
            if not re.match(r"^\d+_", layout.name):
                per_kern[k] = layout
    gekozen, ontbreekt = [], []
    for naam, uitleg in KORTLIJST:
        if naam in per_kern:
            gekozen.append((per_kern[naam], naam, uitleg))
        else:
            ontbreekt.append(naam)
    return gekozen, ontbreekt


def verwijder_slides(prs):
    lijst = prs.slides._sldIdLst
    for sldId in list(lijst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        lijst.remove(sldId)


def snoei_layouts(master, houden):
    """Verwijder elke layout die niet in houden zit."""
    te_houden = {id(l) for l in houden}
    verwijderd = 0
    for layout in list(master.slide_layouts):
        if id(layout) in te_houden:
            continue
        try:
            master.slide_layouts.remove(layout)
            verwijderd += 1
        except Exception as fout:
            print(f"  kon {layout.name} niet verwijderen: {fout}", file=sys.stderr)
    return verwijderd


def knip(tekst, maximum):
    """Kort tekst in tot wat er past, op een woordgrens."""
    if maximum is None or len(tekst) <= maximum:
        return tekst
    afgekapt = tekst[: max(1, maximum - 1)]
    if " " in afgekapt:
        afgekapt = afgekapt.rsplit(" ", 1)[0]
    return afgekapt


def vul_voorbeeld(slide, layout, master, layout_naam, uitleg):
    """Vul de placeholders van een nieuwe slide met neutrale voorbeeldtekst.

    Gebruikt de gemeten rollen en capaciteit, zodat de voorbeeldslide niet overloopt.
    Sjabloonelementen (de 9pt-lijntjes van het bureausjabloon) blijven onaangeroerd.
    """
    kaart = rollen(layout, master)
    for ph in list(slide.placeholders):
        info = kaart.get(ph.placeholder_format.idx, {})
        rol = info.get("rol", "onbekend")
        ruimte = info.get("max_tekens")
        try:
            if rol == "picture":
                zet_beeldplaatshouder(slide, ph)
            elif rol in ("chart", "table"):
                continue
            elif rol == "sjabloonelement":
                continue
            elif rol == "action title":
                ph.text_frame.paragraphs[0].add_run().text = knip(VOORBEELDTEKST["titel"], ruimte)
            elif rol == "tekst":
                ph.text_frame.paragraphs[0].add_run().text = knip(VOORBEELDTEKST["body"], ruimte)
            elif rol == "kort label":
                ph.text_frame.paragraphs[0].add_run().text = knip(VOORBEELDTEKST["cijfer"], ruimte)
            elif rol == "bron":
                ph.text_frame.paragraphs[0].add_run().text = knip(VOORBEELDTEKST["bron"], ruimte)
            elif rol == "extra":
                ph.text_frame.paragraphs[0].add_run().text = knip(VOORBEELDTEKST["body"], ruimte)
        except Exception:
            continue
    notities = slide.notes_slide.notes_text_frame
    notities.text = (
        f"Layout: {layout_naam}. {uitleg}. "
        "Voorbeeldslide in het sjabloon, niet bedoeld om te presenteren."
    )


def zet_beeldplaatshouder(slide, ph):
    """Vervang een beeldplaceholder door een rechthoek met omschrijving."""
    links, boven, breed, hoog = ph.left, ph.top, ph.width, ph.height
    ph._element.getparent().remove(ph._element)
    vorm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, links, boven, breed, hoog)
    vorm.fill.solid()
    vorm.fill.fore_color.rgb = RGBColor(0x9F, 0x9F, 0x9F)
    vorm.line.color.rgb = RGBColor(0x42, 0x42, 0x42)
    vorm.line.width = Pt(1)
    kader = vorm.text_frame
    kader.text = VOORBEELDTEKST["beeld"]
    for para in kader.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x14, 0x14, 0x14)


def main():
    arg = argparse.ArgumentParser(description=__doc__)
    arg.add_argument("bron", help="pad naar het bron-pptx")
    arg.add_argument("-o", "--uit", required=True, help="pad naar het nieuwe sjabloon")
    arg.add_argument("--master", type=int, default=0, help="index van de master (standaard 0)")
    arg.add_argument("--zonder-voorbeelden", action="store_true",
                     help="alleen snoeien, geen voorbeeldslides toevoegen")
    opties = arg.parse_args()

    prs = Presentation(opties.bron)
    master = prs.slide_masters[opties.master]
    print(f"Bron: {opties.bron}")
    print(f"  {len(prs.slides)} slides, {len(prs.slide_masters)} masters, "
          f"{len(master.slide_layouts)} layouts in master {opties.master}")

    gekozen, ontbreekt = kies_layouts(master)
    print(f"  kortlijst: {len(gekozen)} van {len(KORTLIJST)} gevonden")
    if ontbreekt:
        print(f"  NIET gevonden: {', '.join(ontbreekt)}")

    verwijder_slides(prs)
    for extra in list(prs.slide_masters)[1:]:
        print(f"  let op: master '{extra.name}' blijft staan, snoei die apart")

    weg = snoei_layouts(master, [g[0] for g in gekozen])
    print(f"  {weg} layouts verwijderd, {len(master.slide_layouts)} over")

    if not opties.zonder_voorbeelden:
        for layout, naam, uitleg in gekozen:
            slide = prs.slides.add_slide(layout)
            vul_voorbeeld(slide, layout, master, naam, uitleg)
        print(f"  {len(gekozen)} voorbeeldslides toegevoegd")

    prs.save(opties.uit)
    print(f"Weggeschreven: {opties.uit}")


if __name__ == "__main__":
    main()

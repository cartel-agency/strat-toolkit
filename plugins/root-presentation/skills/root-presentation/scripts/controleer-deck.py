#!/usr/bin/env python3
"""Controleer een gebouwd deck op de vormregels van Cartel.

Zes controles, machinaal:
  1. tekstoverloop: staat er meer tekst in een vak dan erin past
  2. gedachtestreepjes, in de slides én in de notities
  3. tekst in hoofdletters alleen
  4. slides zonder spreeknotities
  5. slides met een cijfer maar zonder bronvermelding
  6. lege plaatshouders, opgelijst zodat de strateeg weet wat hij nog moet aanleveren

    python3 controleer-deck.py deck.pptx

Geeft afsluitcode 1 als er fouten zijn, zodat je hem in een keten kan hangen.
"""

import argparse
import os
import re
import sys

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from capaciteit import rollen, TOLERANTIE  # noqa: E402

STREEPJES = ("—", "–")
# Zeven letters of meer in hoofdletters is geen afkorting maar een stijlkeuze,
# en die is bij Cartel niet toegestaan. Korter laten we door: dat zijn afkortingen.
CAPS = re.compile(r"\b[A-ZÀ-Þ]{7,}\b")
CIJFER = re.compile(r"\d+([.,]\d+)?\s*(%|procent|euro|miljoen|miljard)")
BRON = re.compile(r"bron|source|©|volgens ", re.IGNORECASE)
# afkortingen die wel in hoofdletters mogen
CAPS_TOEGELATEN = {"RIZIV", "FAGG", "BTW", "KPI", "KPIS", "ATL", "BTL", "TVDB", "OKR",
                   "SOV", "SOM", "EMAS", "ANSM", "FDA", "ABPI", "CM", "TAM", "MVA"}


def tekst_van(shape):
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text


def capaciteit(shape):
    """Ruwe schatting van wat er in deze vorm past, op basis van de echte maten."""
    if shape.width is None or shape.height is None:
        return None
    grootte = None
    try:
        for para in shape.text_frame.paragraphs:
            if para.font.size:
                grootte = para.font.size.pt
            for run in para.runs:
                if run.font.size:
                    grootte = max(grootte or 0, run.font.size.pt)
    except Exception:
        pass
    if not grootte:
        return None
    breedte_pt = Emu(shape.width).inches * 72.0
    hoogte_pt = Emu(shape.height).inches * 72.0
    per_regel = max(1, int(breedte_pt / (grootte * 0.5)))
    regels = max(1, int(hoogte_pt / (grootte * 1.22)))
    return int(per_regel * regels)


def controleer(pad):
    prs = Presentation(pad)
    fouten, opmerkingen, plaatshouders = [], [], []

    for nummer, slide in enumerate(prs.slides, start=1):
        alle_tekst = []
        heeft_bron = False
        # De layout weet beter hoeveel er past dan de slide zelf: op de slide staat de
        # lettergrootte meestal niet, die erft van de layout.
        try:
            kaart = rollen(slide.slide_layout, prs.slide_masters[0])
        except Exception:
            kaart = {}

        for shape in slide.shapes:
            tekst = tekst_van(shape)
            if not tekst.strip():
                continue
            alle_tekst.append(tekst)

            if "PLAATSHOUDER" in tekst.upper():
                eerste = tekst.split("\n")
                omschrijving = eerste[1] if len(eerste) > 1 else "(geen omschrijving)"
                plaatshouders.append(f"slide {nummer}: {omschrijving}")
                continue

            for streepje in STREEPJES:
                if streepje in tekst:
                    fouten.append(f"slide {nummer}: gedachtestreepje in '{tekst[:60]}'")

            for woord in CAPS.findall(tekst):
                if woord not in CAPS_TOEGELATEN:
                    fouten.append(f"slide {nummer}: '{woord}' staat in hoofdletters")

            ruimte = None
            if shape.is_placeholder:
                info = kaart.get(shape.placeholder_format.idx, {})
                ruimte = info.get("max_tekens")
            if ruimte is None:
                ruimte = capaciteit(shape)
            if ruimte and len(tekst) > ruimte * TOLERANTIE:
                fouten.append(
                    f"slide {nummer}: tekst van {len(tekst)} tekens in een vak voor ongeveer "
                    f"{ruimte}. Loopt waarschijnlijk over: '{tekst[:50]}'"
                )

            if BRON.search(tekst):
                heeft_bron = True

        samen = " ".join(alle_tekst)
        if CIJFER.search(samen) and not heeft_bron:
            opmerkingen.append(f"slide {nummer}: er staat een cijfer op maar geen bron")

        notities = ""
        if slide.has_notes_slide:
            notities = slide.notes_slide.notes_text_frame.text
            for streepje in STREEPJES:
                if streepje in notities:
                    fouten.append(f"slide {nummer}: gedachtestreepje in de spreeknotities")
        if not notities.strip():
            opmerkingen.append(f"slide {nummer}: geen spreeknotities")

    return prs, fouten, opmerkingen, plaatshouders


def main():
    arg = argparse.ArgumentParser(description=__doc__)
    arg.add_argument("deck", help="pad naar het pptx")
    opties = arg.parse_args()

    prs, fouten, opmerkingen, plaatshouders = controleer(opties.deck)
    print(f"{opties.deck}: {len(prs.slides)} slides\n")

    if fouten:
        print(f"FOUTEN ({len(fouten)}), deze moeten weg voor je oplevert:")
        for f in fouten:
            print("  -", f)
    else:
        print("Geen fouten op de vormregels.")

    if opmerkingen:
        print(f"\nAANDACHTSPUNTEN ({len(opmerkingen)}):")
        for o in opmerkingen:
            print("  -", o)

    if plaatshouders:
        print(f"\nPLAATSHOUDERS ({len(plaatshouders)}), dit moet de strateeg nog aanleveren:")
        for p in plaatshouders:
            print("  -", p)

    sys.exit(1 if fouten else 0)


if __name__ == "__main__":
    main()

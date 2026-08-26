#!/usr/bin/env python3
"""Lees een pptx-sjabloon uit en schrijf een layoutinventaris weg als markdown.

Gebruik:
    python3 lees-sjabloon.py sjabloon.pptx > layouts.md
    python3 lees-sjabloon.py sjabloon.pptx --json > layouts.json

Wat het oplevert: per slidelayout de naam, de achtergrondkleur (en dus of het een
donkere of een lichte slide is), en elke placeholder met zijn index, type, positie
en afmeting. Plus het thema: lettertypes en de twaalf themakleuren.

Dat bestand is de keuzelijst waaruit de agent een layout kiest. Hij ontwerpt niet,
hij kiest. Wat hier niet in staat, bestaat voor hem niet.
"""

import argparse
import json
import os
import sys
from collections import Counter

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from capaciteit import rollen  # noqa: E402

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def inch(value):
    if value is None:
        return None
    return round(Emu(value).inches, 2)


def srgb_uit(element):
    """Zoek de eerste expliciete kleurwaarde onder een element."""
    if element is None:
        return None
    srgb = element.find(".//a:srgbClr", NS)
    if srgb is not None:
        return "#" + srgb.get("val", "").upper()
    scheme = element.find(".//a:schemeClr", NS)
    if scheme is not None:
        return "thema:" + scheme.get("val", "")
    return None


def achtergrond_van(layout):
    """Achtergrondkleur van een layout, of de kleurmapping die hij van de master erft."""
    bg = layout.element.find("p:bg", NS)
    kleur = srgb_uit(bg)
    mapping = layout.element.find("p:clrMapOvr/a:overrideClrMapping", NS)
    omgekeerd = False
    if mapping is not None:
        # bg1 die naar dk1 wijst betekent: donkere achtergrond, lichte tekst
        omgekeerd = mapping.get("bg1") in ("dk1", "dk2")
    # De naam is de tweede aanwijzing: niet elke donkere layout draait de mapping om,
    # sommige hebben gewoon een zwarte achtergrondvorm.
    naam = (layout.name or "").lower()
    if "zwart" in naam or "black" in naam:
        omgekeerd = True
    return kleur, omgekeerd


def placeholders_van(layout):
    rijen = []
    for ph in layout.placeholders:
        fmt = ph.placeholder_format
        rijen.append(
            {
                "idx": fmt.idx,
                "type": str(fmt.type).split(" ")[0] if fmt.type is not None else "",
                "naam": ph.name,
                "links": inch(ph.left),
                "boven": inch(ph.top),
                "breedte": inch(ph.width),
                "hoogte": inch(ph.height),
            }
        )
    return sorted(rijen, key=lambda r: (r["boven"] or 0, r["links"] or 0))


def thema_van(prs):
    deel = prs.slide_masters[0].part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    from defusedxml.ElementTree import fromstring

    root = fromstring(deel.blob)
    fonts = {}
    for sleutel, pad in (("titels", "majorFont"), ("tekst", "minorFont")):
        el = root.find(f".//{{{NS['a']}}}{pad}/{{{NS['a']}}}latin")
        fonts[sleutel] = el.get("typeface") if el is not None else None
    kleuren = {}
    schema = root.find(f".//{{{NS['a']}}}clrScheme")
    if schema is not None:
        for kind in schema:
            naam = kind.tag.split("}")[1]
            kleuren[naam] = srgb_uit(kind)
    return {"lettertypes": fonts, "themakleuren": kleuren}


def lees(pad):
    prs = Presentation(pad)
    gebruik = Counter()
    for slide in prs.slides:
        if slide.slide_layout is not None:
            gebruik[slide.slide_layout.name] += 1

    layouts = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            kleur, omgekeerd = achtergrond_van(layout)
            layouts.append(
                {
                    "naam": layout.name,
                    "master": master.name if master.name else "",
                    "achtergrond": kleur,
                    "donkere_kleurmapping": omgekeerd,
                    "gebruikt_in_sjabloon": gebruik.get(layout.name, 0),
                    "placeholders": placeholders_van(layout),
                    "rollen": rollen(layout, master),
                }
            )

    return {
        "bestand": pad,
        "formaat_inch": [inch(prs.slide_width), inch(prs.slide_height)],
        "aantal_slides": len(prs.slides),
        "thema": thema_van(prs),
        "layouts": layouts,
    }


def als_markdown(data):
    uit = []
    b, h = data["formaat_inch"]
    uit.append("# Layoutinventaris van het Cartel-sjabloon\n")
    uit.append(
        "> Automatisch gegenereerd met `scripts/lees-sjabloon.py`. Niet met de hand aanpassen: "
        "verandert het sjabloon, draai het script opnieuw.\n"
    )
    uit.append(f"**Bestand:** `{data['bestand']}`  ")
    uit.append(f"**Formaat:** {b} bij {h} inch  ")
    uit.append(f"**Aantal layouts:** {len(data['layouts'])}\n")

    thema = data["thema"]
    uit.append("## Thema\n")
    uit.append(f"- Titellettertype: `{thema['lettertypes'].get('titels')}`")
    uit.append(f"- Tekstlettertype: `{thema['lettertypes'].get('tekst')}`")
    uit.append("- Themakleuren: " + ", ".join(
        f"`{k}` {v}" for k, v in thema["themakleuren"].items() if v
    ) + "\n")

    uit.append("## Hoe je dit leest\n")
    uit.append(
        "De Cartel-layouts zetten de zichtbare titel **niet** in de titelplaceholder. Die staat op "
        "9pt en draagt het navigatielijntje linksboven. De action title zit in een BODY-placeholder "
        "van 35pt of meer. Vul dus altijd op **rol**, nooit op placeholdertype.\n"
    )
    uit.append(
        "`max` is het gemeten maximum aantal tekens dat in dat vak past. Ga daar niet overheen: "
        "PowerPoint laat tekst gewoon buiten het kader lopen zonder te waarschuwen.\n"
    )
    uit.append(
        "Rollen: `action title` is de conclusiezin, `tekst` de onderbouwing, `kort label` een cijfer "
        "of nummer, `bron` de bronregel, `picture` een beeldvlak (wordt een plaatshouder), "
        "`chart` en `table` blijven leeg tot je ze vult, en `sjabloonelement` raak je nooit aan.\n"
    )

    uit.append("## Layouts\n")
    uit.append("| # | Layout | Donker | Rollen (idx, lettergrootte, max tekens) |")
    uit.append("|---|---|---|---|")
    volgorde = ["action title", "tekst", "kort label", "bron", "extra", "picture", "chart", "table"]
    for i, lay in enumerate(data["layouts"]):
        items = []
        for rol in volgorde:
            for idx, info in sorted(lay["rollen"].items()):
                if info.get("rol") != rol:
                    continue
                if rol in ("picture", "chart", "table"):
                    items.append(f"{rol} `{idx}`")
                else:
                    items.append(f"{rol} `{idx}` ({info.get('pt')}pt, max {info.get('max_tekens')})")
        uit.append(
            f"| {i} | `{lay['naam']}` | {'ja' if lay['donkere_kleurmapping'] else 'nee'} | "
            f"{'; '.join(items) or 'geen invulbare placeholders'} |"
        )
    return "\n".join(uit)


def main():
    arg = argparse.ArgumentParser(description=__doc__)
    arg.add_argument("pptx", help="pad naar het sjabloonbestand")
    arg.add_argument("--json", action="store_true", help="geef json in plaats van markdown")
    opties = arg.parse_args()

    data = lees(opties.pptx)
    if opties.json:
        json.dump(data, sys.stdout, indent=2, ensure_ascii=False)
    else:
        sys.stdout.write(als_markdown(data) + "\n")


if __name__ == "__main__":
    main()
